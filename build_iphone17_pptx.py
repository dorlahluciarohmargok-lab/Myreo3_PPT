#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
iPhone 17 Series — Ultra Premium Apple-Style PPTX
极致苹果风格 · 真实产品图嵌入 · 高级灰渐变 · 25页
"""
import io, os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_MARKER_STYLE

# ─── 苹果高级灰 12阶 ──────────────────────────────────
BK = RGBColor(0x00, 0x00, 0x00)
G1 = RGBColor(0x0A, 0x0A, 0x0A)
G2 = RGBColor(0x11, 0x11, 0x11)
G3 = RGBColor(0x1A, 0x1A, 0x1C)
G4 = RGBColor(0x2C, 0x2C, 0x2E)
G5 = RGBColor(0x3A, 0x3A, 0x3C)
G6 = RGBColor(0x48, 0x48, 0x4A)
G7 = RGBColor(0x63, 0x63, 0x66)
G8 = RGBColor(0x8E, 0x8E, 0x93)
G9 = RGBColor(0xAE, 0xAE, 0xB2)
G10= RGBColor(0xC7, 0xC7, 0xCC)
WH = RGBColor(0xF5, 0xF5, 0xF7)

# ─── 品牌色 · 大胆鲜艳 ──────────────────────────────
BL  = RGBColor(0x00, 0x71, 0xE3)   # Apple Blue
BLL = RGBColor(0x29, 0x97, 0xFF)   # Light Blue
BXL = RGBColor(0x64, 0xD2, 0xFF)   # Extra Light Blue
GR  = RGBColor(0x30, 0xD1, 0x58)   # Green
GRL = RGBColor(0x63, 0xE6, 0xBE)   # Light Green
OR  = RGBColor(0xFF, 0x6B, 0x00)   # Orange
ORL = RGBColor(0xFF, 0x9F, 0x0A)   # Light Orange
PU  = RGBColor(0xBF, 0x5A, 0xF2)   # Purple
PUL = RGBColor(0xDA, 0x8F, 0xFF)   # Light Purple
GO  = RGBColor(0xF7, 0xCE, 0x68)   # Gold
RD  = RGBColor(0xFF, 0x45, 0x3A)   # Red
TE  = RGBColor(0x5A, 0xC8, 0xFA)   # Teal

W = Inches(13.333)
H = Inches(7.5)

IMG_DIR = os.path.join(os.path.dirname(__file__), 'imgs')

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

# ═══ 辅助函数 ══════════════════════════════════════════

def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, c): 
    f = slide.background.fill; f.solid(); f.fore_color.rgb = c

def rect(slide, l, t, w, h, fill, line=None, rnd=False):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if rnd else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(st, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line; s.line.width = Pt(1)
    else: s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h, sz=16, bold=False, col=WH,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = col
    r.font.name = "Helvetica Neue"
    return tb

def txts(slide, lines, l, t, w, h, sz=14, bold=False, col=WH,
         align=PP_ALIGN.LEFT, spacing=1.2):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.size = Pt(sz); r.font.bold = bold
        r.font.color.rgb = col; r.font.name = "Helvetica Neue"
    return tb

def hl(slide, top, col=G4, l=Inches(0.8), w=None):
    w = w or Inches(11.733)
    ln = rect(slide, l, top, w, Pt(1), col)
    return ln

def badge(slide, l, t, label, val, vcol=BLL, cw=Inches(3.5), ch=Inches(1.65)):
    rect(slide, l, t, cw, ch, G3, G5, rnd=True)
    rect(slide, l, t, cw, Pt(4), vcol)          # 顶部彩色线
    txt(slide, label, l+Inches(0.2), t+Inches(0.14), cw-Inches(0.4), Inches(0.4),
        sz=11, col=G8)
    txt(slide, val, l+Inches(0.2), t+Inches(0.58), cw-Inches(0.4), Inches(0.85),
        sz=28, bold=True, col=vcol)

def tag(slide, l, t, text, bgc, tc=WH, tw=Inches(1.8), th=Inches(0.38)):
    rect(slide, l, t, tw, th, bgc, rnd=True)
    txt(slide, text, l, t+Inches(0.04), tw, th-Inches(0.08),
        sz=11, bold=True, col=tc, align=PP_ALIGN.CENTER)

def img(slide, fname, l, t, w, h=None):
    path = os.path.join(IMG_DIR, fname)
    if not os.path.exists(path):
        print(f"  [skip img] {fname}")
        return None
    try:
        if h:
            return slide.shapes.add_picture(path, l, t, w, h)
        return slide.shapes.add_picture(path, l, t, width=w)
    except Exception as e:
        print(f"  [img err] {fname}: {e}")
        return None

def bar_h(slide, l, t, w, h, cats, vals, colors, maxv=None):
    cd = CategoryChartData()
    cd.categories = cats; cd.add_series("", vals)
    cf = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, l, t, w, h, cd)
    ch = cf.chart; ch.has_legend = False
    pl = ch.plots[0]; pl.has_data_labels = True
    pl.data_labels.show_value = True
    pl.data_labels.font.size = Pt(10)
    pl.data_labels.font.color.rgb = WH
    pl.data_labels.font.bold = True
    s = ch.series[0]
    for i, c in enumerate(colors):
        pt = s.points[i]; pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = c
    ch.category_axis.tick_labels.font.color.rgb = G9
    ch.category_axis.tick_labels.font.size = Pt(10)
    ch.value_axis.tick_labels.font.color.rgb = G7
    if maxv: ch.value_axis.maximum_scale = maxv
    return cf

def bar_v(slide, l, t, w, h, cats, vals, colors, maxv=None):
    cd = CategoryChartData()
    cd.categories = cats; cd.add_series("", vals)
    cf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, l, t, w, h, cd)
    ch = cf.chart; ch.has_legend = False
    pl = ch.plots[0]; pl.has_data_labels = True
    pl.data_labels.show_value = True
    pl.data_labels.font.size = Pt(10)
    pl.data_labels.font.color.rgb = WH
    pl.data_labels.font.bold = True
    s = ch.series[0]
    for i, c in enumerate(colors):
        pt = s.points[i]; pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = c
    ch.category_axis.tick_labels.font.color.rgb = G9
    ch.category_axis.tick_labels.font.size = Pt(10)
    ch.value_axis.tick_labels.font.color.rgb = G7
    if maxv: ch.value_axis.maximum_scale = maxv
    return cf

def bar_cluster(slide, l, t, w, h, cats, series):
    cd = CategoryChartData()
    cd.categories = cats
    for name, vals, _ in series: cd.add_series(name, vals)
    cf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, l, t, w, h, cd)
    ch = cf.chart; ch.has_legend = True
    ch.legend.position = XL_LEGEND_POSITION.BOTTOM
    ch.legend.font.size = Pt(11); ch.legend.font.color.rgb = G9
    for i, (_, _, c) in enumerate(series):
        s = ch.series[i]; s.format.fill.solid()
        s.format.fill.fore_color.rgb = c
    ch.category_axis.tick_labels.font.color.rgb = G9
    ch.category_axis.tick_labels.font.size = Pt(10)
    ch.value_axis.tick_labels.font.color.rgb = G7
    return cf

def line_chart(slide, l, t, w, h, cats, series):
    cd = CategoryChartData()
    cd.categories = cats
    for name, vals, _ in series: cd.add_series(name, vals)
    cf = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, l, t, w, h, cd)
    ch = cf.chart; ch.has_legend = True
    ch.legend.position = XL_LEGEND_POSITION.BOTTOM
    ch.legend.font.size = Pt(11); ch.legend.font.color.rgb = G9
    for i, (_, _, c) in enumerate(series):
        s = ch.series[i]
        s.format.line.color.rgb = c; s.format.line.width = Pt(2.5)
        s.marker.style = XL_MARKER_STYLE.CIRCLE; s.marker.size = 7
        s.marker.format.fill.solid(); s.marker.format.fill.fore_color.rgb = c
    ch.category_axis.tick_labels.font.color.rgb = G9
    ch.category_axis.tick_labels.font.size = Pt(10)
    ch.value_axis.tick_labels.font.color.rgb = G7
    return cf

def donut(slide, l, t, w, h, cats, vals, colors):
    cd = CategoryChartData()
    cd.categories = cats; cd.add_series("", vals)
    cf = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, l, t, w, h, cd)
    ch = cf.chart; ch.has_legend = True
    ch.legend.position = XL_LEGEND_POSITION.RIGHT
    ch.legend.font.size = Pt(11); ch.legend.font.color.rgb = G9
    pl = ch.plots[0]; pl.has_data_labels = True
    pl.data_labels.show_percentage = True
    pl.data_labels.font.size = Pt(10); pl.data_labels.font.color.rgb = WH
    s = ch.series[0]
    for i, c in enumerate(colors):
        s.points[i].format.fill.solid()
        s.points[i].format.fill.fore_color.rgb = c
    return cf

print("🚀 Building iPhone 17 Series Ultra Premium PPTX...")
print(f"   图片目录: {IMG_DIR}")
print(f"   可用图片: {len(os.listdir(IMG_DIR))} 张\n")


# ═══════════════════════════════════════════════════
# SLIDE 1 — 封面 Hero（真实产品图 + 深色渐变）
# ═══════════════════════════════════════════════════
print("  [1/25] Hero Cover")
s = blank(); bg(s, BK)
rect(s, Inches(0), Inches(0), W, Pt(4), BL)         # 顶部蓝线
rect(s, Inches(7.2), Inches(0), Inches(6.133), H, G1)  # 右侧深灰区
# 真实产品图 - Pro Max正面图
img(s, 'pro_max_gsm.jpg', Inches(7.3), Inches(0.3), Inches(5.8))
# 左侧内容
tag(s, Inches(0.8), Inches(1.0), "iPhone 17 Series · 2025", BL, WH, Inches(3.0))
txt(s, "Think",      Inches(0.8), Inches(1.8), Inches(6.4), Inches(1.15),
    sz=88, bold=True, col=WH)
txt(s, "Different.", Inches(0.8), Inches(2.85), Inches(6.4), Inches(1.15),
    sz=88, bold=True, col=BLL)
txt(s, "设计之美 · 性能领先 · 影像真实 · 续航卓越",
    Inches(0.8), Inches(4.15), Inches(6.4), Inches(0.55),
    sz=17, col=G8)
# 四款型号标签
for i,(nm,c) in enumerate([("iPhone 17",BLL),("iPhone Air",GR),
                             ("17 Pro",ORL),("Pro Max",GO)]):
    tag(s, Inches(0.8+i*1.62), Inches(5.0), nm, G3, c, Inches(1.55))
# 底部小字
txt(s, "Apple Inc. · Product Presentation 2025 · Confidential",
    Inches(0.8), Inches(7.1), Inches(8), Inches(0.3), sz=9, col=G6)

# ═══════════════════════════════════════════════════
# SLIDE 2 — 产品阵容（真实前视图）
# ═══════════════════════════════════════════════════
print("  [2/25] Lineup Overview")
s = blank(); bg(s, BK)
rect(s, Inches(0), Inches(0), W, Pt(3), G4)
txt(s, "Lineup 2025", Inches(0.7), Inches(0.3), Inches(5), Inches(0.45),
    sz=12, col=G7)
txt(s, "四款 iPhone，一个愿景。", Inches(0.7), Inches(0.75), Inches(11), Inches(0.9),
    sz=50, bold=True, col=WH)
txt(s, "每一款都为特定的你精心打造",
    Inches(0.7), Inches(1.65), Inches(11), Inches(0.45), sz=17, col=G8)
hl(s, Inches(2.2))
# 卡片
card_data = [
    ("iPhone 17",    "从 ¥5,999", ["A19 · 6.3\" ProMotion","双48MP相机系统","Ceramic Shield 2","五色:黑/薰衣草/蓝/绿/白"],    BLL, G3, 'i17_front.jpg'),
    ("iPhone Air",   "从 ¥6,999", ["A19 Pro · 5.64mm超薄","165g钛合金机身","6.5\"大屏超轻","四色:天蓝/星光/黑/紫"],            GR,  G3, 'air_review.jpg'),
    ("iPhone 17 Pro","从 ¥7,999", ["A19 Pro · 三48MP镜头","8×光学品质变焦","铝合金一体机身","橙/蓝/银三色"],                    ORL, G3, 'pro_front.jpg'),
    ("Pro Max",      "从 ¥9,999", ["6.9\" 史上最大屏幕","5,088mAh史上最大电池","全球续航测试第一","256GB～2TB可选"],             GO,  G3, 'pro_max_front.jpg'),
]
for i,(nm,pr,specs,ac,bg_c,im) in enumerate(card_data):
    lf = Inches(0.45 + i*3.22)
    rect(s, lf, Inches(2.4), Inches(3.0), Inches(4.75), bg_c, ac, rnd=True)
    rect(s, lf, Inches(2.4), Inches(3.0), Pt(4), ac)  # 彩色顶线
    # 真实产品图
    img(s, im, lf+Inches(0.1), Inches(2.45), Inches(2.8), Inches(2.0))
    txt(s, nm,  lf+Inches(0.15), Inches(4.55), Inches(2.75), Inches(0.5),
        sz=19, bold=True, col=WH)
    txt(s, pr,  lf+Inches(0.15), Inches(5.05), Inches(2.75), Inches(0.42),
        sz=15, bold=True, col=ac)
    txts(s, specs, lf+Inches(0.15), Inches(5.5), Inches(2.75), Inches(1.4),
         sz=11, col=G9)

# ═══════════════════════════════════════════════════
# SLIDE 3 — 真实外观图：iPhone 17 设计细节
# ═══════════════════════════════════════════════════
print("  [3/25] iPhone 17 Design - Real Photos")
s = blank(); bg(s, G1)
rect(s, Inches(0), Inches(0), W, Pt(4), BL)
rect(s, Inches(6.6), Inches(0), Inches(6.733), H, BK)  # 右暗区
tag(s, Inches(0.7), Inches(0.6), "Design · 外观工艺", BL, WH, Inches(2.0))
txt(s, "精心打磨，",     Inches(0.7), Inches(1.3), Inches(5.8), Inches(0.95),
    sz=56, bold=True, col=WH)
txt(s, "每一处细节。",   Inches(0.7), Inches(2.15), Inches(5.8), Inches(0.95),
    sz=56, bold=True, col=BLL)
txt(s, "轮廓边缘更流线，边框更纤薄。正面 Ceramic Shield 2，抗划性提升3倍。外观之美，经久不衰。",
    Inches(0.7), Inches(3.25), Inches(5.6), Inches(0.85),
    sz=15, col=G9)
# 三项特性卡片
feats = [("🛡️ Ceramic Shield 2","正面抗划性 +3×",BLL),
         ("💧 IP68 防水防尘",   "最深6m · 30分钟",GR),
         ("🔗 MagSafe + USB-C","双模快充·磁吸生态",PUL)]
for i,(ti,de,c) in enumerate(feats):
    tp = Inches(4.35+i*0.85)
    rect(s, Inches(0.7), tp, Inches(5.5), Inches(0.72), G3, G5, rnd=True)
    rect(s, Inches(0.7), tp, Pt(4), Inches(0.72), c)
    txt(s, ti, Inches(1.0), tp+Inches(0.08), Inches(3.5), Inches(0.35), sz=14, bold=True, col=c)
    txt(s, de, Inches(1.0), tp+Inches(0.4), Inches(4.5), Inches(0.27), sz=12, col=G8)
txt(s, "配色：Black · Lavender · Mist Blue · Sage · White",
    Inches(0.7), Inches(7.05), Inches(6), Inches(0.35), sz=12, col=G7)
# 右侧：真实设备图（正面+侧面细节）
img(s, 'hd_002.jpg', Inches(6.8), Inches(0.1), Inches(6.2), Inches(4.8))
img(s, 'lifestyle_001.jpg', Inches(6.8), Inches(5.1), Inches(3.0), Inches(2.25))
img(s, 'hd_003.jpg',       Inches(9.95), Inches(5.1), Inches(3.05), Inches(2.25))

# ═══════════════════════════════════════════════════
# SLIDE 4 — iPhone Air 超薄产品展示
# ═══════════════════════════════════════════════════
print("  [4/25] iPhone Air")
s = blank(); bg(s, BK)
rect(s, Inches(0), Inches(0), W, Pt(4), GR)
rect(s, Inches(6.8), Inches(0), Inches(6.533), H, G1)
tag(s, Inches(0.7), Inches(0.6), "iPhone Air · 2025", GR, BK, Inches(2.0))
txt(s, "史上最薄",      Inches(0.7), Inches(1.35), Inches(5.8), Inches(1.0),
    sz=68, bold=True, col=WH)
txt(s, "iPhone。",      Inches(0.7), Inches(2.25), Inches(5.8), Inches(0.95),
    sz=68, bold=True, col=GR)
txt(s, "仅5.64mm厚，165g轻盈入手。你不敢相信一部iPhone可以这么薄。",
    Inches(0.7), Inches(3.35), Inches(5.8), Inches(0.7), sz=16, col=G9)
badge(s, Inches(0.7),  Inches(4.3), "机身厚度", "5.64mm", GR,  Inches(2.7), Inches(1.6))
badge(s, Inches(3.55), Inches(4.3), "重量",     "165g",   TE,  Inches(2.7), Inches(1.6))
badge(s, Inches(0.7),  Inches(6.0), "屏幕",     "6.5\"",  BLL, Inches(2.7), Inches(1.4))
badge(s, Inches(3.55), Inches(6.0), "芯片",     "A19 Pro", PUL, Inches(2.7), Inches(1.4))
# 右侧：Air 真实产品图
img(s, 'hd_004.jpg', Inches(7.0), Inches(0.1), Inches(6.0), Inches(7.3))

# ═══════════════════════════════════════════════════
# SLIDE 5 — 显示屏参数卡
# ═══════════════════════════════════════════════════
print("  [5/25] Display Specs")
s = blank(); bg(s, BK)
rect(s, Inches(0), Inches(0), W, Pt(4), TE)
tag(s, Inches(0.7), Inches(0.5), "Super Retina XDR", TE, BK, Inches(2.4))
txt(s, "我们最卓越的显示屏。",
    Inches(0.7), Inches(1.05), Inches(11), Inches(0.85),
    sz=52, bold=True, col=WH, align=PP_ALIGN.CENTER)
txt(s, "每一像素精密校准，色彩还原达到专业级精度",
    Inches(0.7), Inches(1.95), Inches(11), Inches(0.45),
    sz=16, col=G8, align=PP_ALIGN.CENTER)
hl(s, Inches(2.5), TE)
specs_disp = [
    ("3000 nits",  "峰值亮度（户外）",  BLL),
    ("120 Hz",     "ProMotion 刷新率",  GR),
    ("460 ppi",    "视网膜像素密度",    ORL),
    ("2,000,000:1","OLED 对比度",       TE),
    ("6.3\"–6.9\"","全系四种尺寸",      PUL),
    ("Always-On",  "Pro 机型常亮显示",  GO),
]
for i,(v,l,c) in enumerate(specs_disp):
    ci = i%3; ri = i//3
    lf = Inches(0.6+ci*4.22); tp = Inches(2.7+ri*2.0)
    rect(s, lf, tp, Inches(3.9), Inches(1.75), G3, G5, rnd=True)
    rect(s, lf, tp, Inches(3.9), Pt(4), c)
    txt(s, v, lf+Inches(0.25), tp+Inches(0.2), Inches(3.4), Inches(0.8),
        sz=32, bold=True, col=c)
    txt(s, l, lf+Inches(0.25), tp+Inches(1.1), Inches(3.4), Inches(0.45),
        sz=13, col=G8)

# ═══════════════════════════════════════════════════
# SLIDE 6 — A19 Pro 芯片（真实芯片图）
# ═══════════════════════════════════════════════════
print("  [6/25] A19 Pro Chip")
s = blank(); bg(s, G1)
rect(s, Inches(0), Inches(0), W, Pt(4), PU)
rect(s, Inches(6.8), Inches(0), Inches(6.533), H, BK)
tag(s, Inches(0.7), Inches(0.55), "A19 Pro · 第二代3nm", PU, WH, Inches(2.6))
txt(s, "领先一切。", Inches(0.7), Inches(1.3), Inches(6), Inches(0.95),
    sz=64, bold=True, col=WH)
txt(s, "比 A18 快15%  ·  比最快Android旗舰快30%+\n神经网络引擎每秒处理 38 万亿次运算",
    Inches(0.7), Inches(2.45), Inches(5.9), Inches(0.9), sz=15, col=G9)
# 性能数据卡
for i,(lb,vl,c) in enumerate([
    ("CPU 提升", "+15%",     PUL),
    ("GPU 提升", "+20%",     BLL),
    ("AI 算力",  "38T ops/s",GR),
]):
    badge(s, Inches(0.7+i*2.1), Inches(3.6), lb, vl, c,
          cw=Inches(1.95), ch=Inches(1.55))
txt(s, "3nm 制程 · 全球移动性能 #1 · 驱动 Apple Intelligence",
    Inches(0.7), Inches(5.4), Inches(6), Inches(0.38), sz=13, col=G7)
hl(s, Inches(5.88), PU, Inches(0.7), Inches(5.5))
# 右侧：真实设备 + 产品细节图
img(s, 'hd_009.jpg',  Inches(7.0), Inches(0.1), Inches(6.0), Inches(3.6))
img(s, 'hd_010.jpg',  Inches(7.0), Inches(3.8), Inches(2.9), Inches(3.55))
img(s, 'hd_011.jpg',  Inches(10.05),Inches(3.8), Inches(2.95), Inches(3.55))


# ═══════════════════════════════════════════════════
# SLIDE 7 — 性能跑分对比（横向条形图）
# ═══════════════════════════════════════════════════
print("  [7/25] Benchmark Chart")
s = blank(); bg(s, BK)
rect(s, Inches(0), Inches(0), W, Pt(4), PU)
txt(s, "性能遥遥领先。", Inches(0.7), Inches(0.4), Inches(11), Inches(0.85),
    sz=48, bold=True, col=WH, align=PP_ALIGN.CENTER)
txt(s, "Geekbench 6 多核跑分 — 分数越高越好",
    Inches(0.7), Inches(1.28), Inches(11), Inches(0.42),
    sz=15, col=G8, align=PP_ALIGN.CENTER)
bar_h(s, Inches(0.8), Inches(1.85), Inches(11.733), Inches(5.35),
      ["iPhone 17 Pro (A19 Pro)", "OnePlus 15 (SD 8 Gen 4)",
       "Samsung S25 Ultra (SD 8 Elite)", "Google Pixel 9 Pro (Tensor G4)"],
      (8950, 7350, 7200, 5800),
      [BLL, G5, G5, G5])

# ═══════════════════════════════════════════════════
# SLIDE 8 — Pro摄像头系统（真实摄像头特写图）
# ═══════════════════════════════════════════════════
print("  [8/25] Pro Camera System - Real Photos")
s = blank(); bg(s, BK)
rect(s, Inches(0), Inches(0), W, Pt(4), OR)
rect(s, Inches(6.5), Inches(0), Inches(6.833), H, G1)
tag(s, Inches(0.7), Inches(0.55), "Pro Camera System", OR, BK, Inches(2.4))
txt(s, "终极 Pro 影像。", Inches(0.7), Inches(1.2), Inches(5.8), Inches(0.9),
    sz=52, bold=True, col=WH)
txt(s, "三颗 48MP Fusion 镜头，等效8支专业镜头。\niPhone史上最长8×光学品质变焦。",
    Inches(0.7), Inches(2.2), Inches(5.8), Inches(0.85), sz=15, col=G9)
cam_items = [
    ("📷  48MP 主摄 Fusion",   "传感器位移防抖 · f/1.78大光圈", OR),
    ("🔭  48MP 长焦 · 8×变焦", "iPhone史上最长光学品质变焦",    ORL),
    ("🌐  48MP 超广角",        "微距 · 170°大视角",              GO),
    ("🎬  4K 120fps ProRes",  "Apple Log 2 · ProRes RAW · genlock", GR),
    ("🤳  18MP Center Stage", "自动取景追踪 · 专业自拍",         BLL),
]
for i,(ti,de,c) in enumerate(cam_items):
    tp = Inches(3.15+i*0.85)
    rect(s, Inches(0.7), tp, Inches(5.5), Inches(0.73), G3, G5, rnd=True)
    rect(s, Inches(0.7), tp, Pt(4), Inches(0.73), c)
    txt(s, ti, Inches(1.0), tp+Inches(0.07), Inches(3.6), Inches(0.37), sz=14, bold=True, col=c)
    txt(s, de, Inches(1.0), tp+Inches(0.42), Inches(4.6), Inches(0.28), sz=12, col=G8)
# 右侧：真实摄像头特写图（最高清、最大图）
img(s, 'hd_014.jpg',  Inches(6.7), Inches(0.05), Inches(6.4), Inches(4.5))  # 主产品图
img(s, 'lifestyle_004.jpg', Inches(6.7),  Inches(4.65), Inches(3.1), Inches(2.7)) # 实拍
img(s, 'lifestyle_005.jpg', Inches(9.95), Inches(4.65), Inches(3.05), Inches(2.7)) # 实拍

# ═══════════════════════════════════════════════════
# SLIDE 9 — Shot on iPhone 实拍展示（大图全页）
# ═══════════════════════════════════════════════════
print("  [9/25] Shot on iPhone - Full Page")
s = blank(); bg(s, BK)
# 左侧文字区
rect(s, Inches(0), Inches(0), Inches(4.8), H, G1)
rect(s, Inches(0), Inches(0), Pt(4), H, OR)  # 左侧橙线
txt(s, "Shot on iPhone", Inches(0.5), Inches(1.2), Inches(4), Inches(0.6),
    sz=13, col=OR, bold=True)
txt(s, "影像，", Inches(0.5), Inches(1.9), Inches(4), Inches(0.9),
    sz=54, bold=True, col=WH)
txt(s, "还原真实。", Inches(0.5), Inches(2.72), Inches(4), Inches(0.9),
    sz=54, bold=True, col=ORL)
txt(s, "每一张照片都是记忆的忠实还原。三颗48MP镜头捕捉世界本来的样子。",
    Inches(0.5), Inches(3.8), Inches(4), Inches(0.85), sz=14, col=G9)
# 数据
for i,(v,lb,c) in enumerate([("#1","DxOMark 排名",OR),
                               ("162","主摄综合分",ORL),
                               ("8×","最长变焦",  GO)]):
    tp = Inches(5.0+i*0.72)
    txt(s, v,  Inches(0.5), tp, Inches(2.5), Inches(0.55), sz=34, bold=True, col=c)
    txt(s, lb, Inches(0.5), tp+Inches(0.45), Inches(2.5), Inches(0.28), sz=11, col=G8)
# 右侧：5张实拍大图网格
img(s, 'hd_017.jpg',  Inches(4.85), Inches(0),    Inches(4.2), Inches(3.75)) # 大图左
img(s, 'hd_018.jpg',  Inches(9.1),  Inches(0),    Inches(4.233),Inches(3.75))# 大图右
img(s, 'hd_019.jpg',  Inches(4.85), Inches(3.76), Inches(2.75), Inches(3.74))# 小图左
img(s, 'hd_020.jpg',  Inches(7.65), Inches(3.76), Inches(2.75), Inches(3.74))# 小图中
img(s, 'hd_021.jpg',  Inches(10.45),Inches(3.76), Inches(2.883),Inches(3.74))# 小图右

# ═══════════════════════════════════════════════════
# SLIDE 10 — 摄像头评分对比（横向条形图）
# ═══════════════════════════════════════════════════
print("  [10/25] Camera Score Chart")
s = blank(); bg(s, BK)
rect(s, Inches(0), Inches(0), W, Pt(4), OR)
txt(s, "影像评分全球第一。", Inches(0.7), Inches(0.4), Inches(11), Inches(0.85),
    sz=48, bold=True, col=WH, align=PP_ALIGN.CENTER)
txt(s, "DxOMark 主摄综合评分 2025年旗舰机型（满分180，越高越好）",
    Inches(0.7), Inches(1.3), Inches(11), Inches(0.42),
    sz=14, col=G8, align=PP_ALIGN.CENTER)
bar_h(s, Inches(0.8), Inches(1.85), Inches(11.733), Inches(5.35),
      ["iPhone 17 Pro Max", "Huawei Pura 70 Ultra",
       "Samsung S25 Ultra", "Google Pixel 9 Pro XL"],
      (162, 151, 155, 148),
      [ORL, G5, G5, G5], maxv=180)

# ═══════════════════════════════════════════════════
# SLIDE 11 — 续航对比（进度条可视化）
# ═══════════════════════════════════════════════════
print("  [11/25] Battery Life")
s = blank(); bg(s, G1)
rect(s, Inches(0), Inches(0), W, Pt(4), GR)
tag(s, Inches(0.7), Inches(0.5), "Battery Life", GR, BK, Inches(1.6))
txt(s, "效率才是真正的续航。",
    Inches(0.7), Inches(1.0), Inches(11), Inches(0.85),
    sz=50, bold=True, col=WH, align=PP_ALIGN.CENTER)
txt(s, "更小的电池，更长的使用时间 — iOS + A19 Pro 协同优化的力量",
    Inches(0.7), Inches(1.9), Inches(11), Inches(0.5),
    sz=16, col=G8, align=PP_ALIGN.CENTER)
hl(s, Inches(2.55), GR)
bat = [
    ("iPhone 17 Pro Max  5,088 mAh", 100, GR,  "🏆 全球 #1 · 3h仅耗9%"),
    ("OnePlus 15  7,300 mAh",         82,  G5,  ""),
    ("Samsung S25 Ultra  5,000 mAh",  78,  G5,  ""),
    ("Google Pixel 9 Pro XL  5,060",  74,  G5,  ""),
]
for i,(nm,pct,c,note) in enumerate(bat):
    tp = Inches(2.75+i*1.05)
    txt(s, nm+(f"  {note}" if note else ""),
        Inches(0.7), tp+Inches(0.08), Inches(4.5), Inches(0.45),
        sz=13, col=WH if c==GR else G9, bold=(c==GR))
    rect(s, Inches(5.4), tp+Inches(0.14), Inches(7.2), Inches(0.38), G4)
    rect(s, Inches(5.4), tp+Inches(0.14), Inches(7.2*pct/100), Inches(0.38), c)
    txt(s, f"{pct}%", Inches(12.75), tp+Inches(0.06), Inches(0.7), Inches(0.45),
        sz=14, bold=True, col=c)
txt(s, "CNET 35款手机横向续航测试 · 3小时全亮度视频流媒体 · iPhone 17 Pro Max 排名第1",
    Inches(0.7), Inches(7.0), Inches(11.5), Inches(0.38), sz=10, col=G6)

# ═══════════════════════════════════════════════════
# SLIDE 12 — 续航数据大数字
# ═══════════════════════════════════════════════════
print("  [12/25] Battery Big Numbers")
s = blank(); bg(s, BK)
rect(s, Inches(0), Inches(0), W, Pt(4), GR)
txt(s, "3小时视频，仅消耗9%电量。",
    Inches(0.7), Inches(0.5), Inches(11), Inches(0.85),
    sz=46, bold=True, col=WH, align=PP_ALIGN.CENTER)
txt(s, "这是 iPhone 17 Pro Max 在 CNET 最严苛测试中的真实表现",
    Inches(0.7), Inches(1.45), Inches(11), Inches(0.45),
    sz=16, col=G8, align=PP_ALIGN.CENTER)
# 大数字三连
for i,(num,lb,desc,c) in enumerate([
    ("#1",       "35款手机续航排名",   "全球综合第一",          GR),
    ("9%",       "3小时视频耗电量",   "CNET测试最优成绩",       TE),
    ("5,088mAh", "Pro Max电池容量",   "iPhone史上最大电池",     ORL),
]):
    lf = Inches(0.8+i*4.2)
    rect(s, lf, Inches(2.3), Inches(3.8), Inches(3.9), G3, c, rnd=True)
    rect(s, lf, Inches(2.3), Inches(3.8), Pt(4), c)
    txt(s, num, lf+Inches(0.2), Inches(2.5), Inches(3.4), Inches(1.4),
        sz=52, bold=True, col=c, align=PP_ALIGN.CENTER)
    txt(s, lb, lf+Inches(0.2), Inches(4.0), Inches(3.4), Inches(0.5),
        sz=14, col=G9, align=PP_ALIGN.CENTER)
    txt(s, desc, lf+Inches(0.2), Inches(4.55), Inches(3.4), Inches(0.45),
        sz=12, col=G7, align=PP_ALIGN.CENTER)
txt(s, "对比：OnePlus 15 搭载 7,300mAh 超大电池，续航测试仍不敌 iPhone 17 Pro Max",
    Inches(0.7), Inches(6.55), Inches(11.5), Inches(0.5),
    sz=14, col=G7, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════
# SLIDE 13 — iOS 26 + Apple Intelligence
# ═══════════════════════════════════════════════════
print("  [13/25] iOS 26 + Apple Intelligence")
s = blank(); bg(s, G1)
rect(s, Inches(0), Inches(0), W, Pt(4), PU)
tag(s, Inches(0.7), Inches(0.5), "iOS 26 · Apple Intelligence", PU, WH, Inches(3.2))
txt(s, "全新外观，更多魔法。",
    Inches(0.7), Inches(1.1), Inches(11), Inches(0.85),
    sz=50, bold=True, col=WH, align=PP_ALIGN.CENTER)
hl(s, Inches(2.1), PU)
feats = [
    ("🧠 Apple Intelligence", "设备端AI·隐私保护",          PUL),
    ("🌐 Live Translation",   "通话/短信实时翻译",           BLL),
    ("🎨 Image Playground",   "一句话生成高质量图像",        ORL),
    ("✍️  Writing Tools",      "智能改写·润色·摘要",          GR),
    ("🎯 Siri 全面进化",       "跨App上下文·真正懂你",        GO),
    ("🔒 Private Cloud",      "联网AI·零数据泄露",           TE),
]
for i,(ti,de,c) in enumerate(feats):
    ci = i%2; ri = i//2
    lf = Inches(0.6+ci*6.4); tp = Inches(2.35+ri*1.7)
    rect(s, lf, tp, Inches(5.9), Inches(1.5), G3, G5, rnd=True)
    rect(s, lf, tp, Pt(4), Inches(1.5), c)
    txt(s, ti, lf+Inches(0.25), tp+Inches(0.15), Inches(5.4), Inches(0.55),
        sz=17, bold=True, col=c)
    txt(s, de, lf+Inches(0.25), tp+Inches(0.75), Inches(5.4), Inches(0.55),
        sz=13, col=G8)

# ═══════════════════════════════════════════════════
# SLIDE 14 — Apple 生态系统
# ═══════════════════════════════════════════════════
print("  [14/25] Apple Ecosystem")
s = blank(); bg(s, BK)
rect(s, Inches(0), Inches(0), W, Pt(4), BL)
tag(s, Inches(0.7), Inches(0.5), "Ecosystem · 生态整合", BL, WH, Inches(2.4))
txt(s, "只有 Apple 才能做到的整合。",
    Inches(0.7), Inches(1.1), Inches(11), Inches(0.85),
    sz=48, bold=True, col=WH, align=PP_ALIGN.CENTER)
txt(s, "iPhone 是生态核心，与每一件 Apple 产品无缝协作",
    Inches(0.7), Inches(2.0), Inches(11), Inches(0.45),
    sz=16, col=G8, align=PP_ALIGN.CENTER)
hl(s, Inches(2.55), BL)
eco = [
    ("⌚ Apple Watch",  "健康数据·无缝解锁",         BLL),
    ("🎧 AirPods",     "自动切换·空间音频",          GR),
    ("💻 Mac",         "Handoff·Universal Control", PUL),
    ("📟 iPad",        "Sidecar副屏·Apple Pencil",  ORL),
    ("📺 Apple TV",    "AirPlay一键投屏·4K无损",    GO),
    ("☁️  iCloud",     "全端同步·15+年安全记录",    TE),
]
for i,(nm,de,c) in enumerate(eco):
    ci = i%3; ri = i//2
    lf = Inches(0.6+ci*4.28); tp = Inches(2.8+ri*2.0)
    rect(s, lf, tp, Inches(3.9), Inches(1.7), G3, G5, rnd=True)
    rect(s, lf, tp, Inches(3.9), Pt(4), c)
    txt(s, nm, lf+Inches(0.2), tp+Inches(0.2), Inches(3.5), Inches(0.6),
        sz=17, bold=True, col=c)
    txt(s, de, lf+Inches(0.2), tp+Inches(0.85), Inches(3.5), Inches(0.6),
        sz=13, col=G8)


# ═══════════════════════════════════════════════════
# SLIDE 15 — iPhone 17 完整规格表
# ═══════════════════════════════════════════════════
print("  [15/25] iPhone 17 Full Specs")
s = blank(); bg(s, BK)
rect(s, Inches(0), Inches(0), W, Pt(4), BLL)
txt(s, "iPhone 17 — 完整规格", Inches(0.7), Inches(0.38), Inches(11), Inches(0.65),
    sz=40, bold=True, col=WH, align=PP_ALIGN.CENTER)
hl(s, Inches(1.15), BLL)
specs17 = [
    ("显示屏", "6.3\" Super Retina XDR OLED · ProMotion 120Hz · 3,000 nits · Dynamic Island"),
    ("芯片",   "Apple A19 · 第二代3nm · 6核CPU + 6核GPU · 第二代神经网络引擎"),
    ("摄像头", "48MP 主摄(OIS) + 48MP 超广角 · 12MP TrueDepth 前摄"),
    ("存储",   "256GB / 512GB"),
    ("电池",   "3,692 mAh · 全天续航 · MagSafe 25W · USB-C 25W"),
    ("颜色",   "Black · Lavender · Mist Blue · Sage · White"),
    ("机身",   "Ceramic Shield 2 正面 · 彩色玻璃背板 · 铝合金边框"),
    ("尺寸",   "149.6 × 71.5 × 7.95mm · 177g"),
    ("防水",   "IP68 · 最深6米，最长30分钟"),
    ("系统",   "iOS 26 · Apple Intelligence · 支持6年系统更新"),
]
for i,(lb,vl) in enumerate(specs17):
    tp = Inches(1.35+i*0.57)
    rect(s, Inches(0.7), tp, Inches(11.733), Inches(0.5),
         G3 if i%2==0 else RGBColor(0x14,0x14,0x16))
    txt(s, lb, Inches(0.9), tp+Inches(0.06), Inches(2.4), Inches(0.42),
        sz=13, col=G8)
    txt(s, vl, Inches(3.5), tp+Inches(0.06), Inches(8.7), Inches(0.42),
        sz=13, col=WH)

# ═══════════════════════════════════════════════════
# SLIDE 16 — iPhone Air 完整规格表
# ═══════════════════════════════════════════════════
print("  [16/25] iPhone Air Full Specs")
s = blank(); bg(s, BK)
rect(s, Inches(0), Inches(0), W, Pt(4), GR)
txt(s, "iPhone Air — 完整规格", Inches(0.7), Inches(0.38), Inches(11), Inches(0.65),
    sz=40, bold=True, col=WH, align=PP_ALIGN.CENTER)
hl(s, Inches(1.15), GR)
specs_air = [
    ("显示屏",   "6.5\" Super Retina XDR OLED · ProMotion 120Hz · Dynamic Island"),
    ("芯片",     "Apple A19 Pro + N1神经芯片 + C1X调制解调器 · 三芯协同"),
    ("厚度",     "5.64mm — iPhone史上最薄"),
    ("重量",     "165g — 比iPhone 6更轻"),
    ("摄像头",   "48MP 主摄 + 12MP TrueDepth 前摄"),
    ("存储",     "128GB / 256GB / 512GB"),
    ("电池",     "全天续航 · USB-C 快充 · MagSafe"),
    ("颜色",     "Sky Blue · Starlight · Space Black · Lilac"),
    ("机身",     "钛金属边框 · Ceramic Shield 正面+背面 · 史上最薄iPhone"),
    ("防水",     "IP68 · 最深6米，最长30分钟"),
]
for i,(lb,vl) in enumerate(specs_air):
    tp = Inches(1.35+i*0.57)
    rect(s, Inches(0.7), tp, Inches(11.733), Inches(0.5),
         G3 if i%2==0 else RGBColor(0x14,0x14,0x16))
    txt(s, lb, Inches(0.9), tp+Inches(0.06), Inches(2.4), Inches(0.42),
        sz=13, col=G8)
    cval = GR if "薄" in vl or "165" in vl else WH
    txt(s, vl, Inches(3.5), tp+Inches(0.06), Inches(8.7), Inches(0.42),
        sz=13, col=cval)

# ═══════════════════════════════════════════════════
# SLIDE 17 — iPhone 17 Pro 完整规格表
# ═══════════════════════════════════════════════════
print("  [17/25] iPhone 17 Pro Full Specs")
s = blank(); bg(s, BK)
rect(s, Inches(0), Inches(0), W, Pt(4), ORL)
txt(s, "iPhone 17 Pro — 完整规格", Inches(0.7), Inches(0.38), Inches(11), Inches(0.65),
    sz=40, bold=True, col=WH, align=PP_ALIGN.CENTER)
hl(s, Inches(1.15), ORL)
specs_pro = [
    ("显示屏",   "6.3\" Super Retina XDR · Always-On · ProMotion 120Hz · 3,000 nits"),
    ("芯片",     "Apple A19 Pro · 第二代3nm · 蒸汽冷却散热"),
    ("主摄",     "48MP Fusion f/1.78 · 传感器位移光学防抖"),
    ("超广角",   "48MP f/2.2 · 微距摄影"),
    ("长焦",     "48MP f/2.8 · 8× 光学品质变焦（iPhone史上最长）"),
    ("前摄",     "18MP Center Stage f/1.9 · 自动取景追踪"),
    ("视频",     "4K 120fps ProRes RAW · Apple Log 2 · genlock · 空间视频"),
    ("存储",     "256GB / 512GB / 1TB / 2TB"),
    ("电池",     "4,685 mAh · 30W USB-C · 25W MagSafe"),
    ("颜色",     "Cosmic Orange · Deep Blue · Silver"),
]
for i,(lb,vl) in enumerate(specs_pro):
    tp = Inches(1.35+i*0.57)
    rect(s, Inches(0.7), tp, Inches(11.733), Inches(0.5),
         G3 if i%2==0 else RGBColor(0x14,0x14,0x16))
    txt(s, lb, Inches(0.9), tp+Inches(0.06), Inches(2.4), Inches(0.42),
        sz=13, col=G8)
    cval = ORL if "变焦" in vl or "最长" in vl else WH
    txt(s, vl, Inches(3.5), tp+Inches(0.06), Inches(8.7), Inches(0.42),
        sz=13, col=cval)

# ═══════════════════════════════════════════════════
# SLIDE 18 — iPhone 17 Pro Max 完整规格表
# ═══════════════════════════════════════════════════
print("  [18/25] iPhone 17 Pro Max Full Specs")
s = blank(); bg(s, BK)
rect(s, Inches(0), Inches(0), W, Pt(4), GO)
txt(s, "iPhone 17 Pro Max — 完整规格", Inches(0.7), Inches(0.38), Inches(11), Inches(0.65),
    sz=40, bold=True, col=WH, align=PP_ALIGN.CENTER)
hl(s, Inches(1.15), GO)
specs_max = [
    ("显示屏",    "6.9\" Super Retina XDR · Always-On · ProMotion 120Hz · 3,000 nits"),
    ("芯片",      "Apple A19 Pro · 第二代3nm · 蒸汽冷却 · 全球性能第一"),
    ("主摄",      "48MP Fusion f/1.78 · 第二代传感器位移防抖"),
    ("超广角",    "48MP f/2.2 · 更高分辨率微距"),
    ("长焦",      "48MP f/2.8 · 8× 光学品质变焦 · 100mm等效焦距"),
    ("前摄",      "18MP Center Stage f/1.9"),
    ("视频",      "4K 120fps ProRes RAW · Apple Log 2 · genlock · 空间视频"),
    ("存储",      "256GB / 512GB / 1TB / 2TB"),
    ("电池",      "5,088 mAh — iPhone史上最大 · 全球续航测试第一"),
    ("颜色",      "Cosmic Orange · Deep Blue · Silver"),
]
for i,(lb,vl) in enumerate(specs_max):
    tp = Inches(1.35+i*0.57)
    rect(s, Inches(0.7), tp, Inches(11.733), Inches(0.5),
         G3 if i%2==0 else RGBColor(0x14,0x14,0x16))
    txt(s, lb, Inches(0.9), tp+Inches(0.06), Inches(2.4), Inches(0.42),
        sz=13, col=G8)
    cval = GO if "史上最大" in vl or "第一" in vl else WH
    txt(s, vl, Inches(3.5), tp+Inches(0.06), Inches(8.7), Inches(0.42),
        sz=13, col=cval)

# ═══════════════════════════════════════════════════
# SLIDE 19 — 全系规格对比矩阵
# ═══════════════════════════════════════════════════
print("  [19/25] Full Comparison Matrix")
s = blank(); bg(s, BK)
rect(s, Inches(0), Inches(0), W, Pt(4), BL)
txt(s, "四款对比，一目了然。", Inches(0.7), Inches(0.3), Inches(11), Inches(0.72),
    sz=44, bold=True, col=WH, align=PP_ALIGN.CENTER)
hl(s, Inches(1.12), BL)
headers = ["", "iPhone 17", "iPhone Air", "17 Pro", "Pro Max"]
hcols   = [None, BLL, GR, ORL, GO]
rows = [
    ("起售价",     "¥5,999", "¥6,999",  "¥7,999",  "¥9,999"),
    ("屏幕",       "6.3\"",  "6.5\"",   "6.3\"",   "6.9\""),
    ("芯片",       "A19",    "A19 Pro",  "A19 Pro", "A19 Pro"),
    ("厚度",       "7.95mm", "5.64mm ✦","8.0mm",   "8.3mm"),
    ("后摄",       "双48MP", "单48MP",  "三48MP",  "三48MP"),
    ("最长变焦",   "2×",     "2×",      "8×",      "8×"),
    ("ProMotion",  "✓",      "✓",       "✓",       "✓"),
    ("Always-On",  "✗",      "✗",       "✓",       "✓"),
    ("ProRes视频", "✗",      "✗",       "✓",       "✓"),
]
# 表头
for ci,hd in enumerate(headers):
    lf = Inches(0.4+ci*2.6)
    if ci>0:
        rect(s, lf, Inches(1.22), Inches(2.45), Inches(0.65), G3, hcols[ci])
    txt(s, hd, lf+Inches(0.1), Inches(1.28), Inches(2.25), Inches(0.52),
        sz=14, bold=True, col=hcols[ci] if hcols[ci] else G7,
        align=PP_ALIGN.CENTER)
# 数据行
for ri,row in enumerate(rows):
    tp = Inches(1.97+ri*0.6)
    fbg = G3 if ri%2==0 else RGBColor(0x15,0x15,0x18)
    for ci,cell in enumerate(row):
        lf = Inches(0.4+ci*2.6)
        rect(s, lf, tp, Inches(2.45), Inches(0.57), fbg if ci>0 else BK)
        col = WH if ci>0 else G7
        if cell=="✓": col = GR
        if cell=="✗": col = G5
        if "✦" in cell: col = GR
        txt(s, cell, lf+Inches(0.1), tp+Inches(0.09), Inches(2.25), Inches(0.42),
            sz=13, bold=(ci==0), col=col,
            align=PP_ALIGN.CENTER if ci>0 else PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════
# SLIDE 20 — 性能趋势折线图（历代比较）
# ═══════════════════════════════════════════════════
print("  [20/25] Performance Trend Chart")
s = blank(); bg(s, BK)
rect(s, Inches(0), Inches(0), W, Pt(4), PU)
txt(s, "代际性能持续跃升。", Inches(0.7), Inches(0.4), Inches(11), Inches(0.85),
    sz=48, bold=True, col=WH, align=PP_ALIGN.CENTER)
txt(s, "Apple 芯片多核跑分历代演进 — Geekbench 6（越高越好）",
    Inches(0.7), Inches(1.3), Inches(11), Inches(0.45),
    sz=15, col=G8, align=PP_ALIGN.CENTER)
line_chart(s, Inches(0.8), Inches(1.85), Inches(11.733), Inches(5.35),
           ["A15", "A16", "A17 Pro", "A18", "A18 Pro", "A19", "A19 Pro"],
           [
               ("iPhone CPU", (5200, 5900, 7100, 7800, 8200, 8400, 8950), BLL),
               ("Android 最强", (4800, 5500, 6100, 6700, 7000, 7200, 7350), G6),
           ])

# ═══════════════════════════════════════════════════
# SLIDE 21 — 续航折线图（历代对比）
# ═══════════════════════════════════════════════════
print("  [21/25] Battery History Chart")
s = blank(); bg(s, BK)
rect(s, Inches(0), Inches(0), W, Pt(4), GR)
txt(s, "续航持续突破。", Inches(0.7), Inches(0.4), Inches(11), Inches(0.85),
    sz=48, bold=True, col=WH, align=PP_ALIGN.CENTER)
txt(s, "Pro Max 视频播放续航时长（小时）历代演进",
    Inches(0.7), Inches(1.3), Inches(11), Inches(0.45),
    sz=15, col=G8, align=PP_ALIGN.CENTER)
line_chart(s, Inches(0.8), Inches(1.85), Inches(11.733), Inches(5.35),
           ["iPhone 13 Pro Max","iPhone 14 Pro Max","iPhone 15 Pro Max",
            "iPhone 16 Pro Max","iPhone 17 Pro Max"],
           [("视频播放续航(h)", (28, 29, 29, 33, 36), GRL)])

# ═══════════════════════════════════════════════════
# SLIDE 22 — 真实产品实拍细节（摄像头+背部+侧面）
# ═══════════════════════════════════════════════════
print("  [22/25] Real Product Detail Photos")
s = blank(); bg(s, BK)
rect(s, Inches(0), Inches(0), W, Pt(4), G4)
txt(s, "精密工艺，每一处都值得细看。",
    Inches(0.7), Inches(0.35), Inches(11), Inches(0.75),
    sz=44, bold=True, col=WH, align=PP_ALIGN.CENTER)
# 6张产品细节图大图布局
img(s, 'hd_023.jpg',       Inches(0.1),  Inches(1.2),  Inches(4.35), Inches(3.3))  # 摄像头特写
img(s, 'hd_028.jpg',       Inches(4.55), Inches(1.2),  Inches(4.35), Inches(3.3))  # 机身细节
img(s, 'hd_036.jpg',       Inches(8.95), Inches(1.2),  Inches(4.233),Inches(3.3))  # 侧面/背部
img(s, 'lifestyle_013.jpg',Inches(0.1),  Inches(4.65), Inches(4.35), Inches(2.75)) # 实拍1
img(s, 'lifestyle_014.jpg',Inches(4.55), Inches(4.65), Inches(4.35), Inches(2.75)) # 实拍2
img(s, 'lifestyle_015.jpg',Inches(8.95), Inches(4.65), Inches(4.233),Inches(2.75)) # 实拍3
# 图注
for i,(lb,x) in enumerate([
    ("摄像头模组特写", Inches(0.1)),
    ("铝合金一体机身", Inches(4.55)),
    ("背部细节工艺",   Inches(8.95)),
]):
    txt(s, lb, x+Inches(0.1), Inches(4.45), Inches(4.1), Inches(0.35),
        sz=11, col=G7)

# ═══════════════════════════════════════════════════
# SLIDE 23 — 六大购买理由
# ═══════════════════════════════════════════════════
print("  [23/25] 6 Reasons to Buy")
s = blank(); bg(s, G1)
rect(s, Inches(0), Inches(0), W, Pt(4), BL)
txt(s, "为什么选择 iPhone 17 系列？",
    Inches(0.7), Inches(0.4), Inches(11), Inches(0.8),
    sz=46, bold=True, col=WH, align=PP_ALIGN.CENTER)
txt(s, "六个理由，让你下定决心",
    Inches(0.7), Inches(1.25), Inches(11), Inches(0.45),
    sz=16, col=G8, align=PP_ALIGN.CENTER)
reasons = [
    ("🎨 外观质感",  "每一道曲线精密计算·Ceramic Shield2·手感奢华",   BLL),
    ("⚡ 性能领先",  "A19 Pro全球移动端第一·30%+超越Android·卡顿绝迹",ORL),
    ("📸 影像真实",  "三颗48MP · 8×变焦 · DxOMark #1 · 4K120fps ProRes",GR),
    ("🔋 续航更好",  "5,088mAh+iOS优化 · CNET 35机第一 · 3h仅耗9%",    GO),
    ("🌐 生态闭环",  "iPhone+Mac+iPad+Watch · 无缝协作 · 生产力翻倍",    PUL),
    ("🛡️ 系统丝滑",  "iOS 26滑动如丝 · 6年系统更新 · Apple Intelligence",TE),
]
for i,(ti,de,c) in enumerate(reasons):
    ci = i%2; ri = i//2
    lf = Inches(0.6+ci*6.4); tp = Inches(1.9+ri*1.75)
    rect(s, lf, tp, Inches(5.9), Inches(1.55), G3, G5, rnd=True)
    rect(s, lf, tp, Pt(4), Inches(1.55), c)
    txt(s, ti, lf+Inches(0.25), tp+Inches(0.15), Inches(5.4), Inches(0.55),
        sz=17, bold=True, col=c)
    txt(s, de, lf+Inches(0.25), tp+Inches(0.78), Inches(5.4), Inches(0.6),
        sz=12, col=G9)

# ═══════════════════════════════════════════════════
# SLIDE 24 — 购买页（定价 + CTA）
# ═══════════════════════════════════════════════════
print("  [24/25] Buy Page")
s = blank(); bg(s, BK)
rect(s, Inches(0), Inches(0), W, Pt(4), BL)
txt(s, "今天，就选你的下一部 iPhone。",
    Inches(0.7), Inches(0.4), Inches(11), Inches(0.8),
    sz=46, bold=True, col=WH, align=PP_ALIGN.CENTER)
hl(s, Inches(1.3), BL)
buy_data = [
    ("iPhone 17",    "从 ¥5,999", "256GB起 · 5色", BLL, 'i17_front.jpg'),
    ("iPhone Air",   "从 ¥6,999", "128GB起 · 史上最薄", GR, 'air_review.jpg'),
    ("iPhone 17 Pro","从 ¥7,999", "256GB起 · Pro影像", ORL, 'pro_front.jpg'),
    ("Pro Max",      "从 ¥9,999", "256GB~2TB · 续航第一", GO, 'pro_max_gsm.jpg'),
]
for i,(nm,pr,desc,c,im) in enumerate(buy_data):
    lf = Inches(0.45+i*3.22)
    rect(s, lf, Inches(1.5), Inches(3.0), Inches(5.6), G3, c, rnd=True)
    rect(s, lf, Inches(1.5), Inches(3.0), Pt(4), c)
    # 产品图
    img(s, im, lf+Inches(0.1), Inches(1.55), Inches(2.8), Inches(2.8))
    # 文字
    txt(s, nm, lf+Inches(0.15), Inches(4.5), Inches(2.75), Inches(0.5),
        sz=18, bold=True, col=WH)
    txt(s, pr, lf+Inches(0.15), Inches(4.98), Inches(2.75), Inches(0.45),
        sz=22, bold=True, col=c)
    txt(s, desc, lf+Inches(0.15), Inches(5.45), Inches(2.75), Inches(0.35),
        sz=12, col=G8)
    # 购买按钮（矩形模拟）
    rect(s, lf+Inches(0.4), Inches(6.0), Inches(2.2), Inches(0.55), c, rnd=True)
    txt(s, "立即购买 →", lf+Inches(0.4), Inches(6.05), Inches(2.2), Inches(0.45),
        sz=13, bold=True, col=BK if c==GR or c==GO else WH,
        align=PP_ALIGN.CENTER)
txt(s, "apple.com/cn  ·  以旧换新  ·  分期免息  ·  教育优惠",
    Inches(0.7), Inches(7.1), Inches(11.5), Inches(0.32),
    sz=12, col=G6, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════
# SLIDE 25 — 结尾（真实产品大图）
# ═══════════════════════════════════════════════════
print("  [25/25] Closing Slide")
s = blank(); bg(s, BK)
rect(s, Inches(0), Inches(0), W, Pt(4), BL)
# 左侧内容
txt(s, "Apple · 2025", Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.45),
    sz=13, col=G7, bold=False)
txt(s, "Think",      Inches(0.8), Inches(1.95), Inches(5.5), Inches(1.2),
    sz=82, bold=True, col=WH)
txt(s, "Different.", Inches(0.8), Inches(3.05), Inches(5.5), Inches(1.2),
    sz=82, bold=True, col=BLL)
txt(s, "iPhone 17 系列 · 重新定义每一种可能。",
    Inches(0.8), Inches(4.4), Inches(5.5), Inches(0.55), sz=16, col=G8)
for i,(nm,c) in enumerate([("外观",BLL),("续航#1",GR),("影像",ORL),("芯片",PUL),("生态",TE)]):
    tag(s, Inches(0.8+i*1.12), Inches(5.2), nm, G3, c, Inches(1.0))
txt(s, "apple.com", Inches(0.8), Inches(6.45), Inches(3), Inches(0.45),
    sz=15, col=BLL, bold=True)
# 右侧：真实大图
img(s, 'hd_014.jpg', Inches(6.3), Inches(0.1), Inches(7.0), Inches(5.0))
img(s, 'hd_023.jpg', Inches(6.3), Inches(5.15), Inches(3.4), Inches(2.25))
img(s, 'hd_036.jpg', Inches(9.85),Inches(5.15), Inches(3.3), Inches(2.25))

# ═══════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════
OUT = "/projects/sandbox/Myreo3_PPT/Apple_iPhone17_Series.pptx"
prs.save(OUT)
print(f"\n✅  Done! → {OUT}")
print(f"   Total slides: {len(prs.slides)}")
print(f"   File size: {os.path.getsize(OUT)//1024}KB")
