#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
iPhone 17 Series - Apple-Style PPTX Generator
25 slides, full Apple aesthetic, product images from apple.com CDN
"""
import urllib.request
import io
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData

# ─── APPLE COLOR SYSTEM ─────────────────────────────────────────────────────
BLACK   = RGBColor(0x00, 0x00, 0x00)
WHITE   = RGBColor(0xF5, 0xF5, 0xF7)
GRAY    = RGBColor(0x6E, 0x6E, 0x73)
LGRAY   = RGBColor(0xD2, 0xD2, 0xD7)
BLUE    = RGBColor(0x06, 0xC0, 0xFF)
GREEN   = RGBColor(0x30, 0xD1, 0x58)
PURPLE  = RGBColor(0xBF, 0x5A, 0xF2)
ORANGE  = RGBColor(0xFF, 0x9F, 0x0A)
GOLD    = RGBColor(0xF7, 0xCE, 0x68)
RED     = RGBColor(0xFF, 0x45, 0x3A)
DGRAY   = RGBColor(0x1C, 0x1C, 0x1E)
MGRAY   = RGBColor(0x2C, 0x2C, 0x2E)

# ─── PRODUCT IMAGE URLS (Apple CDN / newsroom) ───────────────────────────────
IMGS = {
    "hero17":    "https://www.apple.com/newsroom/images/product/iphone/standard/Apple_iPhone-17_Design_09092025_Full-Bleed-Image.jpg.small.jpg",
    "hero_air":  "https://www.apple.com/newsroom/images/product/iphone/standard/Apple_iPhone-Air_Design_09092025_Full-Bleed-Image.jpg.small.jpg",
    "hero_pro":  "https://www.apple.com/newsroom/images/product/iphone/standard/Apple_iPhone-17-Pro_Design_09092025_Full-Bleed-Image.jpg.small.jpg",
    "camera17":  "https://www.apple.com/newsroom/images/product/iphone/standard/Apple_iPhone-17_Camera_09092025_Full-Bleed-Image.jpg.small.jpg",
    "cam_pro":   "https://www.apple.com/newsroom/images/product/iphone/standard/Apple_iPhone-17-Pro_Camera_09092025_Full-Bleed-Image.jpg.small.jpg",
    "chip":      "https://www.apple.com/newsroom/images/product/iphone/standard/Apple_iPhone-17-Pro_Chip_09092025_Full-Bleed-Image.jpg.small.jpg",
    "display":   "https://www.apple.com/newsroom/images/product/iphone/standard/Apple_iPhone-17_Display_09092025_Full-Bleed-Image.jpg.small.jpg",
    "lifestyle": "https://www.apple.com/newsroom/images/product/iphone/standard/Apple_iPhone-17-Pro_Lifestyle_09092025_Full-Bleed-Image.jpg.small.jpg",
}

W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

# ─── HELPERS ────────────────────────────────────────────────────────────────
def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color=BLACK):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

def txt(slide, text, left, top, w, h, size=18, bold=False,
        color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Helvetica Neue"
    p.alignment = align
    return tb

def ctxt(slide, text, top, h=Inches(1), size=48, bold=True, color=WHITE, italic=False):
    return txt(slide, text, Inches(0.5), top, Inches(12.333), h,
               size, bold, color, PP_ALIGN.CENTER, italic)

def sub(slide, text, top, size=20, color=GRAY):
    return txt(slide, text, Inches(1), top, Inches(11.333), Inches(0.7),
               size, False, color, PP_ALIGN.CENTER)

def rect(slide, left, top, w, h, fill=DGRAY, line=None, radius=True):
    shp_type = 5 if radius else 1  # rounded rect vs rect
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius
                                   else MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape

def pill(slide, left, top, w, h, fill, text, tsize=13, tcol=WHITE):
    r = rect(slide, left, top, w, h, fill=fill)
    txt(slide, text, left, top + Inches(0.05), w, h, tsize, True, tcol, PP_ALIGN.CENTER)

def badge(slide, left, top, label, value, vcol=WHITE, lsize=11, vsize=28):
    rect(slide, left, top, Inches(3.8), Inches(1.8), fill=MGRAY)
    txt(slide, label, left+Inches(0.25), top+Inches(0.25),
        Inches(3.3), Inches(0.5), lsize, False, GRAY)
    txt(slide, value, left+Inches(0.25), top+Inches(0.75),
        Inches(3.3), Inches(0.8), vsize, True, vcol)

def fetch_img(url):
    try:
        req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0"})
        data = urllib.request.urlopen(req, timeout=8).read()
        return io.BytesIO(data)
    except Exception:
        return None

def add_img(slide, url, left, top, w, h=None):
    bio = fetch_img(url)
    if bio:
        try:
            if h:
                pic = slide.shapes.add_picture(bio, left, top, w, h)
            else:
                pic = slide.shapes.add_picture(bio, left, top, width=w)
            return pic
        except Exception:
            pass
    return None

def hline(slide, top, color=MGRAY):
    from pptx.util import Pt as UPt
    line = slide.shapes.add_shape(1, Inches(1), top, Inches(11.333), Pt(1))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()

def spec_row(slide, label, value, top, vcol=WHITE):
    txt(slide, label, Inches(1.5), top, Inches(4), Inches(0.5), 14, False, GRAY)
    txt(slide, value, Inches(6),   top, Inches(6), Inches(0.5), 14, True,  vcol)
    hline(slide, top + Inches(0.45), MGRAY)

def bar_chart(slide, left, top, w, h, labels, vals, colors):
    from pptx.chart.data import CategoryChartData
    cd = CategoryChartData()
    cd.categories = labels
    cd.add_series("", vals)
    cf = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, left, top, w, h, cd)
    chart = cf.chart
    chart.has_legend = False
    plot  = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.show_value = True
    plot.data_labels.font.color.rgb = WHITE
    plot.data_labels.font.size = Pt(11)
    s = chart.series[0]
    for i, c in enumerate(colors):
        s.points[i].format.fill.solid()
        s.points[i].format.fill.fore_color.rgb = c
    chart.category_axis.tick_labels.font.color.rgb = WHITE
    chart.category_axis.tick_labels.font.size = Pt(11)
    chart.value_axis.tick_labels.font.color.rgb = GRAY
    chart.value_axis.tick_labels.font.size = Pt(10)
    return cf



# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — HERO COVER
# ═══════════════════════════════════════════════════════════════════════════
s = blank(); bg(s, BLACK)
add_img(s, IMGS["hero17"], Inches(0), Inches(0), W, H)
# dark overlay bottom half
ov = s.shapes.add_shape(1, Inches(0), Inches(4.2), W, Inches(3.3))
ov.fill.solid(); ov.fill.fore_color.rgb = BLACK
ov.line.fill.background()
from pptx.util import Pt as UPt
from pptx.oxml.ns import qn
import lxml.etree as etree
txt(s, "\uf8ff", Inches(6.2), Inches(1.0), Inches(1), Inches(0.9), 42, False, WHITE, PP_ALIGN.CENTER)
txt(s, "iPhone 17 Series", Inches(1), Inches(4.5), Inches(11.333), Inches(1.4),
    72, True, WHITE, PP_ALIGN.CENTER)
txt(s, "Design. Power. Brilliance.", Inches(1), Inches(5.8), Inches(11.333), Inches(0.7),
    28, False, GRAY, PP_ALIGN.CENTER)
txt(s, "iPhone 17  ·  iPhone Air  ·  iPhone 17 Pro  ·  iPhone 17 Pro Max",
    Inches(1), Inches(6.5), Inches(11.333), Inches(0.6),
    16, False, LGRAY, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — LINEUP OVERVIEW
# ═══════════════════════════════════════════════════════════════════════════
s = blank(); bg(s, BLACK)
ctxt(s, "Four iPhones. One Vision.", Inches(0.5), Inches(0.9), 44, True, WHITE)
sub(s, "Each one precisely made for who you are and how you live.", Inches(1.5))

cards = [
    ("iPhone 17",      "从 $799",  "A19 · 6.3\" · 48MP双摄",   BLUE,   "Black / Lavender / Mist Blue / Sage / White"),
    ("iPhone Air",     "从 $899",  "A19 Pro · 5.64mm · 超薄",   GREEN,  "Sky Blue / Starlight / Black"),
    ("iPhone 17 Pro",  "从 $999",  "A19 Pro · 6.3\" · 48MP三摄", ORANGE, "Cosmic Orange / Deep Blue / Silver"),
    ("iPhone 17 Pro Max","从 $1199","A19 Pro · 6.9\" · 全球续航第一", GOLD, "Cosmic Orange / Deep Blue / Silver"),
]
for i, (name, price, spec, col, colors) in enumerate(cards):
    lf = Inches(0.4 + i * 3.2)
    rect(s, lf, Inches(2.4), Inches(3.0), Inches(4.6), fill=DGRAY)
    pill(s, lf + Inches(0.2), Inches(2.6), Inches(1.4), Inches(0.38), col, "NEW", 11)
    txt(s, name,  lf+Inches(0.2), Inches(3.1), Inches(2.7), Inches(0.6), 20, True,  WHITE)
    txt(s, price, lf+Inches(0.2), Inches(3.65), Inches(2.7), Inches(0.5), 16, True,  col)
    txt(s, spec,  lf+Inches(0.2), Inches(4.15), Inches(2.7), Inches(0.7), 12, False, LGRAY)
    txt(s, colors,lf+Inches(0.2), Inches(4.85), Inches(2.7), Inches(0.7), 10, False, GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — DESIGN: iPhone 17
# ═══════════════════════════════════════════════════════════════════════════
s = blank(); bg(s, BLACK)
add_img(s, IMGS["hero17"], Inches(6.8), Inches(0.3), Inches(5.8))
txt(s, "Design", Inches(0.8), Inches(1.2), Inches(5.5), Inches(0.6), 14, False, BLUE)
txt(s, "精心打磨\n每一处细节。", Inches(0.8), Inches(1.8), Inches(5.8), Inches(2.0),
    48, True, WHITE)
txt(s, "轮廓边缘更流线，边框更纤薄，正面搭载 Ceramic Shield 2 玻璃——抗划性能提升 3 倍。\n外观之美，经久不衰。",
    Inches(0.8), Inches(3.9), Inches(5.8), Inches(1.6), 18, False, LGRAY)
for i, (c, name) in enumerate([("#1C1C1E","Black"),("#E2E0EF","Lavender"),("#BDD5E7","Mist Blue"),("#9BB89A","Sage"),("#F5F5F7","White")]):
    cx = Inches(0.8 + i * 0.72)
    dot = s.shapes.add_shape(9, cx, Inches(5.8), Inches(0.45), Inches(0.45))
    dot.fill.solid()
    r,g,b = int(c[1:3],16), int(c[3:5],16), int(c[5:7],16)
    dot.fill.fore_color.rgb = RGBColor(r,g,b)
    dot.line.color.rgb = MGRAY
    txt(s, name, cx - Inches(0.1), Inches(6.3), Inches(0.9), Inches(0.4), 9, False, GRAY, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — DESIGN: iPhone Air (超薄)
# ═══════════════════════════════════════════════════════════════════════════
s = blank(); bg(s, BLACK)
add_img(s, IMGS["hero_air"], Inches(6.5), Inches(0), Inches(6.8))
txt(s, "iPhone Air", Inches(0.8), Inches(1.2), Inches(5.5), Inches(0.6), 14, False, GREEN)
txt(s, "史上最薄\niPhone。", Inches(0.8), Inches(1.8), Inches(5.5), Inches(1.8), 52, True, WHITE)
txt(s, "仅 5.64mm 厚，165g 轻盈入手。\n超越 iPhone 6 成为有史以来最薄 iPhone。\nCeramic Shield 正反两面，比以往任何 iPhone 更耐用。",
    Inches(0.8), Inches(3.8), Inches(5.5), Inches(1.8), 18, False, LGRAY)
for val, lbl, col in [("5.64mm","机身厚度",GREEN),("165g","重量",WHITE),("6.5\"","屏幕尺寸",BLUE)]:
    idx = [("5.64mm","机身厚度",GREEN),("165g","重量",WHITE),("6.5\"","屏幕尺寸",BLUE)].index((val,lbl,col))
    badge(s, Inches(0.8 + idx*3.0), Inches(5.6), lbl, val, col)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — DISPLAY
# ═══════════════════════════════════════════════════════════════════════════
s = blank(); bg(s, BLACK)
add_img(s, IMGS["display"], Inches(0), Inches(0), W, H)
# simple dark overlay for readability
ov2 = s.shapes.add_shape(1, Inches(0), Inches(0), W, Inches(2.8))
ov2.fill.solid(); ov2.fill.fore_color.rgb = BLACK; ov2.line.fill.background()
ctxt(s, "Super Retina XDR", Inches(0.4), Inches(0.5), 52, True, WHITE)
ctxt(s, "我们有史以来最卓越的显示屏。", Inches(0.8), Inches(1.4), 24, False, GRAY)
specs_disp = [
    ("6.3\" / 6.5\" / 6.3\" / 6.9\"", "全系显示尺寸", WHITE),
    ("3000 nits",                       "峰值亮度（户外）", BLUE),
    ("ProMotion 120Hz",                 "自适应刷新率", GREEN),
    ("2,000,000:1",                     "对比度", LGRAY),
    ("460 ppi",                         "像素密度", ORANGE),
    ("Always-On Display",               "Pro 机型常亮屏", PURPLE),
]
for i, (val, lbl, col) in enumerate(specs_disp):
    col_i = i % 3; row_i = i // 3
    lf = Inches(1.0 + col_i * 3.8); tp = Inches(3.2 + row_i * 1.7)
    rect(s, lf, tp, Inches(3.4), Inches(1.4), fill=RGBColor(0x14,0x14,0x18))
    txt(s, val, lf+Inches(0.2), tp+Inches(0.15), Inches(3.0), Inches(0.6), 20, True, col)
    txt(s, lbl, lf+Inches(0.2), tp+Inches(0.75), Inches(3.0), Inches(0.45), 12, False, GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — A19 CHIP PERFORMANCE
# ═══════════════════════════════════════════════════════════════════════════
s = blank(); bg(s, BLACK)
add_img(s, IMGS["chip"], Inches(7.0), Inches(0.5), Inches(5.5))
txt(s, "芯片", Inches(0.8), Inches(1.0), Inches(5.5), Inches(0.5), 14, False, PURPLE)
txt(s, "A19 Pro\n领先一切。", Inches(0.8), Inches(1.6), Inches(5.8), Inches(2.0), 52, True, WHITE)
txt(s, "第二代 3nm 工艺打造。比上一代快 15%，\n比市面最快 Android 旗舰快 30%+。\n最新一代神经网络引擎，每秒 38 万亿次运算，\n驱动 Apple Intelligence 全部功能。",
    Inches(0.8), Inches(3.7), Inches(5.8), Inches(2.0), 17, False, LGRAY)
for i, (label, val, col) in enumerate([
    ("CPU 提升", "+15% vs A18", GREEN),
    ("GPU 提升", "+20% vs A18", BLUE),
    ("AI 算力",  "38T ops/s",   PURPLE),
]):
    badge(s, Inches(0.8 + i*3.0), Inches(5.85), label, val, col, 11, 22)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — A19 vs ANDROID BENCHMARK (BAR CHART)
# ═══════════════════════════════════════════════════════════════════════════
s = blank(); bg(s, BLACK)
ctxt(s, "性能遥遥领先。", Inches(0.4), Inches(0.9), 44, True, WHITE)
sub(s, "Geekbench 6 多核跑分对比（越高越好）", Inches(1.4))
labels = ["iPhone 17 Pro\nA19 Pro", "Samsung S25 Ultra\nSnapdragon 8 Elite", "Google Pixel 9 Pro\nTensor G4", "OnePlus 15\nSD 8 Gen 4"]
vals   = (8950, 7200, 5800, 7350)
colors = [BLUE, LGRAY, LGRAY, LGRAY]
bar_chart(s, Inches(1.5), Inches(2.2), Inches(10.5), Inches(4.8), labels, vals, colors)
txt(s, "* 数据来源：Geekbench 官方数据库 2025年数据", Inches(1), Inches(7.1), Inches(10), Inches(0.3), 10, False, GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — iPHONE 17 PRO DESIGN
# ═══════════════════════════════════════════════════════════════════════════
s = blank(); bg(s, BLACK)
add_img(s, IMGS["hero_pro"], Inches(6.5), Inches(0), Inches(6.8))
txt(s, "iPhone 17 Pro", Inches(0.8), Inches(1.2), Inches(5.5), Inches(0.6), 14, False, ORANGE)
txt(s, "航天级铝合金\n一体机身。", Inches(0.8), Inches(1.8), Inches(5.5), Inches(2.0), 48, True, WHITE)
txt(s, "热锻铝合金一体结构——更坚固、更轻盈、散热更优秀。\n告别钛合金接缝，一体成形天然优雅。\n每一个角度都是工业设计的杰作。",
    Inches(0.8), Inches(3.9), Inches(5.5), Inches(1.6), 18, False, LGRAY)
for i, (c, nm) in enumerate([(0xFF6B2C,"Cosmic Orange"),(0x1D3A5E,"Deep Blue"),(0xC0C0C0,"Silver")]):
    cx = Inches(0.8 + i*1.6)
    d = s.shapes.add_shape(9, cx, Inches(5.8), Inches(0.52), Inches(0.52))
    d.fill.solid(); d.fill.fore_color.rgb = RGBColor((c>>16)&0xFF,(c>>8)&0xFF,c&0xFF)
    d.line.color.rgb = LGRAY
    txt(s, nm, cx - Inches(0.1), Inches(6.35), Inches(1.8), Inches(0.45), 9, False, GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — CAMERA SYSTEM: iPhone 17 Pro
# ═══════════════════════════════════════════════════════════════════════════
s = blank(); bg(s, BLACK)
add_img(s, IMGS["cam_pro"], Inches(0), Inches(0), W, H)
ov3 = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(6.5), H)
ov3.fill.solid(); ov3.fill.fore_color.rgb = BLACK; ov3.line.fill.background()
txt(s, "影像", Inches(0.8), Inches(1.0), Inches(5.5), Inches(0.5), 14, False, ORANGE)
txt(s, "终极\nPro 相机系统。", Inches(0.8), Inches(1.6), Inches(5.8), Inches(2.0), 48, True, WHITE)
feats = [
    ("三颗 48MP Fusion 镜头", "主摄 · 超广角 · 长焦，等效 8 支镜头"),
    ("8× 光学品质变焦",       "iPhone 史上最长变焦"),
    ("18MP Center Stage 前摄","自拍直接跃升专业级"),
    ("4K 120fps ProRes 视频", "媲美专业摄影机的电影录制"),
]
for i, (title, desc) in enumerate(feats):
    tp = Inches(4.0 + i*0.78)
    txt(s, "▸ " + title, Inches(0.8), tp, Inches(5.6), Inches(0.38), 15, True, WHITE)
    txt(s, desc,         Inches(1.2), tp+Inches(0.32), Inches(5.3), Inches(0.35), 13, False, GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — CAMERA: iPhone 17 (普通版)
# ═══════════════════════════════════════════════════════════════════════════
s = blank(); bg(s, BLACK)
add_img(s, IMGS["camera17"], Inches(6.8), Inches(0.5), Inches(5.8))
txt(s, "Camera", Inches(0.8), Inches(1.0), Inches(5.5), Inches(0.5), 14, False, BLUE)
txt(s, "48MP 双摄\n超震撼画质。", Inches(0.8), Inches(1.6), Inches(5.5), Inches(1.8), 48, True, WHITE)
txt(s, "全新 48MP 主摄 + 48MP 超广角，\n配合下一代传感器位移式光学防抖，\n无论白昼夜晚都能拍出令人屏息的画面。",
    Inches(0.8), Inches(3.6), Inches(5.5), Inches(1.5), 18, False, LGRAY)
for i, (lbl, val, col) in enumerate([
    ("主摄像素", "48MP", WHITE), ("超广角", "48MP", BLUE), ("前置", "12MP TrueDepth", GREEN)
]):
    rect(s, Inches(0.8+i*3.0), Inches(5.5), Inches(2.7), Inches(1.55), fill=MGRAY)
    txt(s, lbl, Inches(1.0+i*3.0), Inches(5.65), Inches(2.3), Inches(0.4), 11, False, GRAY)
    txt(s, val, Inches(1.0+i*3.0), Inches(6.05), Inches(2.3), Inches(0.6), 18, True,  col)



# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — CAMERA COMPARISON (PIE/DONUT style via bar)
# ═══════════════════════════════════════════════════════════════════════════
s = blank(); bg(s, BLACK)
ctxt(s, "影像，还原真实。", Inches(0.4), Inches(0.8), 44, True, WHITE)
sub(s, "DxOMark 主摄评分对比（2025年旗舰机型）", Inches(1.4))
lbs = ["iPhone 17 Pro Max", "Samsung S25 Ultra", "Google Pixel 9 Pro XL", "Huawei Pura 70 Ultra"]
vs  = (162, 155, 148, 151)
cs  = [ORANGE, LGRAY, LGRAY, LGRAY]
bar_chart(s, Inches(1.5), Inches(2.2), Inches(10.5), Inches(4.5), lbs, vs, cs)
txt(s, "DxOMark 评分越高代表整体影像表现越优秀，数据截至 2025 年 Q4", Inches(1), Inches(7.0), Inches(11), Inches(0.35), 10, False, GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — BATTERY LIFE (续航)
# ═══════════════════════════════════════════════════════════════════════════
s = blank(); bg(s, BLACK)
txt(s, "续航", Inches(0.8), Inches(0.8), Inches(5), Inches(0.5), 14, False, GREEN)
ctxt(s, "效率才是真正的续航。", Inches(0.4), Inches(1.2), 44, True, WHITE)
sub(s, "更小的电池，更长的使用时间。这就是芯片 + 系统协同优化的力量。", Inches(1.9))
# Battery comparison bars
bat_data = [
    ("iPhone 17 Pro Max\n5,088 mAh", 100, GREEN),
    ("OnePlus 15\n7,300 mAh", 82, LGRAY),
    ("Samsung S25 Ultra\n5,000 mAh", 78, LGRAY),
    ("Google Pixel 9 Pro XL\n5,060 mAh", 74, LGRAY),
]
for i,(lbl,pct,col) in enumerate(bat_data):
    tp = Inches(2.8 + i * 1.0)
    txt(s, lbl, Inches(0.5), tp+Inches(0.08), Inches(3.5), Inches(0.7), 13, False, LGRAY if col==LGRAY else WHITE)
    track = s.shapes.add_shape(1, Inches(4.0), tp+Inches(0.18), Inches(8.5), Inches(0.38))
    track.fill.solid(); track.fill.fore_color.rgb = MGRAY; track.line.fill.background()
    fill_w = Inches(8.5 * pct / 100)
    bar = s.shapes.add_shape(5, Inches(4.0), tp+Inches(0.18), fill_w, Inches(0.38))
    bar.fill.solid(); bar.fill.fore_color.rgb = col; bar.line.fill.background()
    txt(s, f"{pct}%", Inches(12.7), tp+Inches(0.1), Inches(0.6), Inches(0.45), 13, True, col)
txt(s, "基准：CNET 35 机型横向续航测试，3h 视频流媒体，满亮度  |  iPhone 17 Pro Max 排名第一",
    Inches(0.5), Inches(7.0), Inches(12.5), Inches(0.35), 10, False, GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13 — PRO MAX BATTERY DETAIL
# ═══════════════════════════════════════════════════════════════════════════
s = blank(); bg(s, BLACK)
ctxt(s, "3 小时视频\n仅消耗 9% 电量。", Inches(0.5), Inches(1.0), 56, True, WHITE)
sub(s, "这是 iPhone 17 Pro Max 在 CNET 最严苛测试中的真实表现", Inches(2.8))
# Big number highlight
for i, (num, lbl, col) in enumerate([
    ("#1", "35 款手机中续航排名第一", GREEN),
    ("9%", "3h 视频消耗电量", BLUE),
    ("5,088mAh", "Pro Max 电池容量（iPhone 史上最大）", ORANGE),
]):
    lf = Inches(0.8 + i*4.1)
    rect(s, lf, Inches(3.5), Inches(3.7), Inches(3.0), fill=MGRAY)
    txt(s, num, lf+Inches(0.25), Inches(3.7), Inches(3.2), Inches(1.2), 48, True, col, PP_ALIGN.CENTER)
    txt(s, lbl, lf+Inches(0.25), Inches(4.9), Inches(3.2), Inches(0.9), 13, False, LGRAY, PP_ALIGN.CENTER)
txt(s, "对比：OnePlus 15 搭载 7,300mAh 超大电池，续航测试中仍不敌 iPhone 17 Pro Max",
    Inches(1), Inches(6.8), Inches(11.333), Inches(0.5), 14, False, GRAY, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 14 — iOS 26 & APPLE INTELLIGENCE
# ═══════════════════════════════════════════════════════════════════════════
s = blank(); bg(s, DGRAY)
ctxt(s, "iOS 26", Inches(0.4), Inches(0.6), 64, True, WHITE)
sub(s, "全新外观。更多魔法。Apple Intelligence 赋能。", Inches(1.4))
feats_ios = [
    ("Apple Intelligence", "基于设备端 AI，保护隐私，理解个人上下文", PURPLE, "🧠"),
    ("Live Translation",   "实时翻译通话、短信和 FaceTime，打破语言壁垒", BLUE,   "🌐"),
    ("Image Playground",   "一句话生成高质量图像，无限创意随心释放", ORANGE, "🎨"),
    ("Writing Tools",      "智能改写、润色、摘要，让表达更精准", GREEN,  "✍️"),
    ("Siri 全面进化",       "深度整合 App 操作，跨 App 理解上下文意图", GOLD,   "🎯"),
    ("隐私保护",            "所有 AI 推理优先在本地运行，敏感数据不上云", WHITE,  "🔒"),
]
for i, (title, desc, col, icon) in enumerate(feats_ios):
    ci = i % 2; ri = i // 2
    lf = Inches(0.8 + ci*6.2); tp = Inches(2.2 + ri*1.6)
    rect(s, lf, tp, Inches(5.6), Inches(1.35), fill=MGRAY)
    txt(s, icon+" "+title, lf+Inches(0.25), tp+Inches(0.15), Inches(5.0), Inches(0.5), 16, True, col)
    txt(s, desc,           lf+Inches(0.25), tp+Inches(0.65), Inches(5.0), Inches(0.55), 13, False, LGRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 15 — ECOSYSTEM (生态)
# ═══════════════════════════════════════════════════════════════════════════
s = blank(); bg(s, BLACK)
txt(s, "Ecosystem", Inches(0.8), Inches(0.8), Inches(5), Inches(0.5), 14, False, BLUE)
ctxt(s, "只有 Apple\n才能做到的整合。", Inches(0.4), Inches(1.2), 44, True, WHITE)
sub(s, "iPhone 是 Apple 生态的核心，与每一件 Apple 产品无缝协作。", Inches(2.2))
eco_items = [
    ("Apple Watch",   "健康数据自动同步，无缝解锁", "⌚", BLUE),
    ("AirPods",       "自动切换，空间音频随行", "🎧", GREEN),
    ("Mac",           "Handoff / Universal Control 跨设备工作流", "💻", PURPLE),
    ("iPad",          "Sidecar 副屏，Apple Pencil 配合", "📱", ORANGE),
    ("Apple TV",      "AirPlay 一键投屏，4K 画质无损", "📺", RED),
    ("iCloud",        "照片、文件、备份全端同步", "☁️", LGRAY),
]
for i, (dev, desc, icon, col) in enumerate(eco_items):
    ci = i % 3; ri = i // 2
    lf = Inches(0.8 + ci*4.2); tp = Inches(3.0 + ri*1.9)
    rect(s, lf, tp, Inches(3.8), Inches(1.55), fill=MGRAY)
    txt(s, icon+" "+dev, lf+Inches(0.2), tp+Inches(0.18), Inches(3.4), Inches(0.5), 17, True, col)
    txt(s, desc,         lf+Inches(0.2), tp+Inches(0.72), Inches(3.4), Inches(0.6), 13, False, LGRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 16 — CONTINUITY & AIRDROP
# ═══════════════════════════════════════════════════════════════════════════
s = blank(); bg(s, DGRAY)
ctxt(s, "设备之间，\n无缝流动。", Inches(0.5), Inches(0.8), 52, True, WHITE)
sub(s, "iPhone 与 Mac、iPad、Apple Watch 之间的协作，流畅到让你忘记切换的存在。", Inches(2.4))
conts = [
    ("AirDrop", "极速传输，附近设备秒收文件、照片、链接", "📡", GREEN),
    ("Handoff", "在 iPhone 开始的任务，拿起 Mac 立刻继续", "🔁", BLUE),
    ("iPhone 镜像", "在 Mac 上直接操控 iPhone，无需拿起手机", "🖥️", ORANGE),
    ("短信/通话转发", "iPhone 上的电话短信，Mac 和 iPad 同步接听", "📞", PURPLE),
]
for i, (title, desc, icon, col) in enumerate(conts):
    lf = Inches(0.8 + (i%2)*6.2); tp = Inches(3.0 + (i//2)*2.0)
    rect(s, lf, tp, Inches(5.6), Inches(1.6), fill=MGRAY)
    txt(s, icon+" "+title, lf+Inches(0.25), tp+Inches(0.18), Inches(5.0), Inches(0.5), 18, True, col)
    txt(s, desc,           lf+Inches(0.25), tp+Inches(0.75), Inches(5.0), Inches(0.6), 14, False, LGRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 17 — PRIVACY & SECURITY
# ═══════════════════════════════════════════════════════════════════════════
s = blank(); bg(s, BLACK)
ctxt(s, "隐私。\n是我们的根本。", Inches(0.4), Inches(0.8), 52, True, WHITE)
sub(s, "从芯片到软件，每一层都是你的保护盾。", Inches(2.3))
priv_items = [
    ("Face ID", "神经网络人脸识别，误识率低于百万分之一", WHITE),
    ("本地 AI 推理", "Apple Intelligence 优先在设备本地运行，数据不离机", GREEN),
    ("Private Cloud Compute", "需联网时，服务器无法访问或保留您的个人数据", BLUE),
    ("App 跟踪透明度", "每个 App 必须征得您同意才能追踪行为", ORANGE),
    ("邮件隐私保护", "阻止发件方追踪您是否已打开邮件", PURPLE),
    ("安全检查", "一键审核并撤销过往授予他人的所有权限", RED),
]
for i, (title, desc, col) in enumerate(priv_items):
    ci = i%2; ri = i//2
    lf = Inches(0.7 + ci*6.3); tp = Inches(2.6 + ri*1.55)
    rect(s, lf, tp, Inches(5.8), Inches(1.2), fill=DGRAY)
    txt(s, "🔒 "+title, lf+Inches(0.25), tp+Inches(0.1), Inches(5.2), Inches(0.45), 15, True, col)
    txt(s, desc,        lf+Inches(0.25), tp+Inches(0.6), Inches(5.2), Inches(0.5),  13, False, LGRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 18 — iPHONE 17 FULL SPECS TABLE
# ═══════════════════════════════════════════════════════════════════════════
s = blank(); bg(s, BLACK)
ctxt(s, "iPhone 17 — 完整规格", Inches(0.4), Inches(0.4), 36, True, WHITE)
hline(s, Inches(1.15))
specs = [
    ("显示屏", "6.3\" Super Retina XDR OLED · ProMotion 120Hz · 3000nits"),
    ("芯片",   "Apple A19 · 第二代 3nm · 6核CPU + 6核GPU"),
    ("摄像头", "48MP 主摄(OIS) + 48MP 超广角 · 12MP TrueDepth 前摄"),
    ("存储",   "256GB / 512GB"),
    ("电池",   "3,692 mAh · 全天续航 · MagSafe + USB-C"),
    ("颜色",   "Black · Lavender · Mist Blue · Sage · White"),
    ("机身",   "Ceramic Shield 2 正面 · 彩色玻璃背板 · 铝合金边框"),
    ("尺寸",   "149.6 × 71.5 × 7.95mm · 177g"),
    ("防水",   "IP68 · 最深 6 米，最长 30 分钟"),
    ("系统",   "iOS 26 · Apple Intelligence · 支持 6 年系统更新"),
]
for i,(label,value) in enumerate(specs):
    spec_row(s, label, value, Inches(1.3 + i*0.53))

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 19 — iPHONE AIR FULL SPECS TABLE
# ═══════════════════════════════════════════════════════════════════════════
s = blank(); bg(s, BLACK)
ctxt(s, "iPhone Air — 完整规格", Inches(0.4), Inches(0.4), 36, True, WHITE)
hline(s, Inches(1.15))
specs_air = [
    ("显示屏",  "6.5\" Super Retina XDR OLED · ProMotion 120Hz · 2000nits"),
    ("芯片",    "Apple A19 Pro · C1X 调制解调器芯片 · N1 神经处理器"),
    ("机身厚度", "5.64mm（iPhone 史上最薄）"),
    ("重量",    "165g"),
    ("摄像头",  "48MP 主摄 · 12MP TrueDepth 前摄"),
    ("存储",    "128GB / 256GB / 512GB"),
    ("电池",    "全天续航 · USB-C 快充 · MagSafe"),
    ("颜色",    "Sky Blue · Starlight · Black"),
    ("机身",    "Ceramic Shield 正面+背面 · 铝合金框架"),
    ("防水",    "IP68 · 最深 6 米，最长 30 分钟"),
]
for i,(label,value) in enumerate(specs_air):
    spec_row(s, label, value, Inches(1.3 + i*0.53), GREEN if "薄" in value or "165" in value else WHITE)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 20 — iPHONE 17 PRO / PRO MAX FULL SPECS
# ═══════════════════════════════════════════════════════════════════════════
s = blank(); bg(s, BLACK)
ctxt(s, "iPhone 17 Pro / Pro Max — 完整规格", Inches(0.4), Inches(0.4), 32, True, WHITE)
hline(s, Inches(1.1))
specs_pro = [
    ("显示屏",    "Pro: 6.3\" · Pro Max: 6.9\" · Always-On · ProMotion 120Hz · 3000nits"),
    ("芯片",      "Apple A19 Pro · 第二代 3nm · 业界最强移动芯片"),
    ("摄像头",    "三颗 48MP Fusion（主摄+超广角+长焦）· 8× 光学品质变焦"),
    ("前置",      "18MP Center Stage · 自动取景追踪"),
    ("视频",      "4K 120fps ProRes · Log 录制 · Apple Log"),
    ("存储",      "256GB / 512GB / 1TB / 2TB"),
    ("电池",      "Pro Max: 5,088 mAh（史上最大）· 全球 35 机型续航测试第一"),
    ("机身",      "热锻铝合金一体机身 · Ceramic Shield 正面"),
    ("颜色",      "Cosmic Orange · Deep Blue · Silver"),
    ("尺寸/重量", "Pro: 149.6mm · Pro Max: 163mm · 227g"),
]
for i,(label,value) in enumerate(specs_pro):
    spec_row(s, label, value, Inches(1.2 + i*0.53), ORANGE if "续航" in value or "5,088" in value else WHITE)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 21 — SIDE-BY-SIDE COMPARISON (全系对比)
# ═══════════════════════════════════════════════════════════════════════════
s = blank(); bg(s, BLACK)
ctxt(s, "四款对比，一目了然。", Inches(0.4), Inches(0.4), 38, True, WHITE)
headers = ["", "iPhone 17", "iPhone Air", "iPhone 17 Pro", "iPhone 17 Pro Max"]
rows = [
    ("起售价",    "$799",    "$899",   "$999",   "$1199"),
    ("屏幕",      "6.3\"",   "6.5\"",  "6.3\"",  "6.9\""),
    ("芯片",      "A19",     "A19 Pro","A19 Pro", "A19 Pro"),
    ("后摄",      "双48MP",  "单48MP", "三48MP", "三48MP"),
    ("最长变焦",  "2×",      "2×",     "8×",     "8×"),
    ("厚度",      "7.95mm",  "5.64mm", "8.0mm",  "8.3mm"),
    ("ProMotion", "✓",       "✓",      "✓",      "✓"),
    ("Always-On", "✗",       "✗",      "✓",      "✓"),
    ("ProRes 视频","✗",      "✗",      "✓",      "✓"),
]
col_colors = [None, BLUE, GREEN, ORANGE, GOLD]
for ci, hdr in enumerate(headers):
    lf = Inches(0.3 + ci*2.55)
    rect(s, lf, Inches(1.1), Inches(2.45), Inches(0.6), fill=MGRAY if ci>0 else RGBColor(0,0,0))
    txt(s, hdr, lf+Inches(0.15), Inches(1.15), Inches(2.2), Inches(0.5),
        14, True, col_colors[ci] if col_colors[ci] else GRAY, PP_ALIGN.CENTER)
for ri, row in enumerate(rows):
    tp = Inches(1.8 + ri*0.58)
    fill_row = MGRAY if ri%2==0 else RGBColor(0x1A,0x1A,0x1C)
    for ci, cell in enumerate(row):
        lf = Inches(0.3 + ci*2.55)
        rect(s, lf, tp, Inches(2.45), Inches(0.55), fill=fill_row if ci>0 else RGBColor(0,0,0))
        col = WHITE if ci>0 else GRAY
        if cell == "✓": col = GREEN
        if cell == "✗": col = RGBColor(0x44,0x44,0x44)
        txt(s, cell, lf+Inches(0.15), tp+Inches(0.08), Inches(2.2), Inches(0.42),
            13, ci==0, col, PP_ALIGN.CENTER if ci>0 else PP_ALIGN.LEFT)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 22 — SHOT ON iPHONE (影像展示)
# ═══════════════════════════════════════════════════════════════════════════
s = blank(); bg(s, BLACK)
add_img(s, IMGS["lifestyle"], Inches(0), Inches(0), W, H)
o4 = s.shapes.add_shape(1, Inches(0), Inches(5.5), W, Inches(2.0))
o4.fill.solid(); o4.fill.fore_color.rgb = BLACK; o4.line.fill.background()
txt(s, "Shot on iPhone 17 Pro Max", Inches(0.8), Inches(5.7), Inches(12), Inches(0.6),
    32, True, WHITE, PP_ALIGN.LEFT)
txt(s, "每一张照片都是真实的还原，每一个瞬间都值得被记住。",
    Inches(0.8), Inches(6.25), Inches(10), Inches(0.55), 18, False, LGRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 23 — WHY iPHONE (购买理由)
# ═══════════════════════════════════════════════════════════════════════════
s = blank(); bg(s, DGRAY)
ctxt(s, "为什么选择 iPhone？", Inches(0.4), Inches(0.5), 44, True, WHITE)
sub(s, "六个理由，让你下定决心。", Inches(1.3))
reasons = [
    ("🎨 外观质感", "每一道曲线都经过精密计算，拿在手中就是奢华", BLUE),
    ("⚡ 性能领先",  "A19 Pro 跑分碾压所有 Android 旗舰，游戏、AI 毫无卡顿", ORANGE),
    ("📸 影像真实",  "三颗 48MP 镜头，8× 变焦，DxOMark 全球第一", GREEN),
    ("🔋 续航更好",  "5,088mAh + iOS 优化，35 机型测试全球第一", GOLD),
    ("🌐 生态闭环",  "iPhone + Mac + iPad + Watch，无缝协作，生产力翻倍", PURPLE),
    ("🛡️ 系统丝滑",  "iOS 26 每次滑动都是享受，6 年系统更新保值", WHITE),
]
for i,(title,desc,col) in enumerate(reasons):
    ci = i%2; ri = i//2
    lf = Inches(0.8 + ci*6.2); tp = Inches(2.0 + ri*1.65)
    rect(s, lf, tp, Inches(5.7), Inches(1.4), fill=MGRAY)
    txt(s, title, lf+Inches(0.25), tp+Inches(0.12), Inches(5.2), Inches(0.5), 18, True, col)
    txt(s, desc,  lf+Inches(0.25), tp+Inches(0.68), Inches(5.2), Inches(0.55), 13, False, LGRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 24 — PRICING & CTA
# ═══════════════════════════════════════════════════════════════════════════
s = blank(); bg(s, BLACK)
ctxt(s, "今天就体验\n你的下一部 iPhone。", Inches(0.5), Inches(0.8), 48, True, WHITE)
sub(s, "所有机型现已开放选购", Inches(2.6))
pricing = [
    ("iPhone 17",       "从 $799 起",   "256GB 起", BLUE),
    ("iPhone Air",      "从 $899 起",   "128GB 起", GREEN),
    ("iPhone 17 Pro",   "从 $999 起",   "256GB 起", ORANGE),
    ("iPhone 17 Pro Max","从 $1199 起", "256GB 起", GOLD),
]
for i,(name,price,storage,col) in enumerate(pricing):
    lf = Inches(0.5 + i*3.2)
    rect(s, lf, Inches(3.2), Inches(3.0), Inches(3.2), fill=MGRAY)
    pill(s, lf+Inches(0.2), Inches(3.4), Inches(2.0), Inches(0.38), col, name, 11)
    txt(s, price,   lf+Inches(0.2), Inches(3.95), Inches(2.7), Inches(0.7), 26, True,  col)
    txt(s, storage, lf+Inches(0.2), Inches(4.65), Inches(2.7), Inches(0.45), 13, False, GRAY)
    pill(s, lf+Inches(0.45), Inches(5.5), Inches(2.1), Inches(0.48), col, "立即购买 →", 13)

txt(s, "apple.com/store  ·  支持以旧换新  ·  分期付款零利率",
    Inches(1), Inches(6.9), Inches(11.333), Inches(0.45), 14, False, GRAY, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 25 — CLOSING
# ═══════════════════════════════════════════════════════════════════════════
s = blank(); bg(s, BLACK)
add_img(s, IMGS["hero_pro"], Inches(3.5), Inches(0.3), Inches(6.5))
o5 = s.shapes.add_shape(1, Inches(0), Inches(5.0), W, Inches(2.5))
o5.fill.solid(); o5.fill.fore_color.rgb = BLACK; o5.line.fill.background()
txt(s, "\uf8ff", Inches(6.1), Inches(5.2), Inches(1.2), Inches(0.9), 44, False, WHITE, PP_ALIGN.CENTER)
ctxt(s, "Think Different.", Inches(5.4), Inches(1.5), 52, True, WHITE)
sub(s, "iPhone 17 系列  ·  2025 年 9 月上市", Inches(6.2), 16, GRAY)
txt(s, "apple.com", Inches(5.5), Inches(6.9), Inches(2.333), Inches(0.45), 14, False, BLUE, PP_ALIGN.CENTER)

# ─── SAVE ────────────────────────────────────────────────────────────────
OUT = "Apple_iPhone17_Series.pptx"
prs.save(OUT)
print(f"✅  Saved → {OUT}  ({len(prs.slides)} slides)")
